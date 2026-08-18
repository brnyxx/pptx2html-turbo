from __future__ import annotations

import io
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


def _text(shape: object, value: str, size: int, color: str = "172033") -> None:
    frame = shape.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = value
    run.font.size = Pt(size)
    run.font.color.rgb = RGBColor.from_string(color)


def _box(
    slide: object,
    left: float,
    top: float,
    width: float,
    height: float,
    value: str,
    fill: str,
    *,
    size: int = 16,
    shape_type: MSO_SHAPE = MSO_SHAPE.ROUNDED_RECTANGLE,
    text_color: str = "172033",
) -> object:
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill)
    shape.line.color.rgb = RGBColor.from_string("D7DCE5")
    _text(shape, value, size, text_color)
    return shape


def _title(slide: object, value: str, subtitle: str) -> None:
    title = slide.shapes.add_textbox(
        Inches(0.55), Inches(0.25), Inches(8.7), Inches(0.65)
    )
    frame = title.text_frame
    run = frame.paragraphs[0].add_run()
    run.text = value
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string("172033")
    sub = slide.shapes.add_textbox(
        Inches(9.35), Inches(0.35), Inches(3.4), Inches(0.4)
    )
    _text(sub, subtitle, 10, "667085")


def _sample_image(label: str, first: str, second: str) -> bytes:
    image = Image.new("RGB", (640, 360), f"#{first}")
    draw = ImageDraw.Draw(image)
    for y in range(360):
        ratio = y / 359
        rgb = tuple(
            round(
                int(first[index : index + 2], 16) * (1 - ratio)
                + int(second[index : index + 2], 16) * ratio
            )
            for index in (0, 2, 4)
        )
        draw.line((0, y, 640, y), fill=rgb)
    draw.ellipse(
        (220, 70, 420, 270),
        fill=(255, 255, 255),
        outline=(23, 32, 51),
        width=8,
    )
    draw.text((264, 164), label, fill=(23, 32, 51))
    stream = io.BytesIO()
    image.save(stream, format="PNG")
    return stream.getvalue()


def _dashboard(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title(slide, "Operations Dashboard", "SLIDE 1 / EXECUTIVE VIEW")
    cards = [
        ("Revenue", "$4.82M", "DDF4E8"),
        ("Conversion", "18.6%", "E5EDFF"),
        ("Risk", "Low", "FFF0D8"),
    ]
    for index, (label, value, color) in enumerate(cards):
        card = _box(
            slide,
            0.6 + index * 4.15,
            1.2,
            3.75,
            1.35,
            f"{label}\n{value}",
            color,
            size=20,
        )
        card.shadow.inherit = False
    _box(slide, 0.6, 2.9, 8.0, 3.8, "Quarterly throughput", "F8FAFC", size=14)
    for index, height in enumerate([1.1, 1.8, 1.4, 2.6, 2.2, 3.0, 2.55]):
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.1 + index * 0.95),
            Inches(6.15 - height),
            Inches(0.55),
            Inches(height),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor.from_string("4F7DF3")
        bar.line.fill.background()
    _box(
        slide,
        8.9,
        2.9,
        3.8,
        3.8,
        "Milestones\n\n✓ Discovery\n✓ Prototype\n→ Validation\n○ Release",
        "172033",
        size=15,
    )


def _text_layout(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title(slide, "Typography & Flow", "SLIDE 2 / TEXT STRESS")
    for index, heading in enumerate(["Strategy", "Execution", "Signals"]):
        panel = _box(
            slide, 0.6 + index * 4.15, 1.25, 3.75, 5.7, "", "F8FAFC"
        )
        frame = panel.text_frame
        frame.clear()
        frame.margin_left = Inches(0.25)
        frame.margin_right = Inches(0.25)
        first = frame.paragraphs[0]
        run = first.add_run()
        run.text = heading
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string("172033")
        items = [
            (0, "One measurable outcome"),
            (1, "Supporting evidence"),
            (1, "Owner and deadline"),
            (
                0,
                "Multiline wrapping validates narrow-column metrics and paragraph spacing.",
            ),
        ]
        for level, value in items:
            paragraph = frame.add_paragraph()
            paragraph.level = level
            paragraph.text = value
            paragraph.font.size = Pt(14 if level == 0 else 12)
            paragraph.font.color.rgb = RGBColor.from_string("172033")
            paragraph.space_before = Pt(10)


def _shapes(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title(slide, "Shape Laboratory", "SLIDE 3 / GEOMETRY")
    kinds = [
        MSO_SHAPE.CHEVRON,
        MSO_SHAPE.HEXAGON,
        MSO_SHAPE.STAR_5_POINT,
        MSO_SHAPE.ARC,
        MSO_SHAPE.RIGHT_ARROW,
        MSO_SHAPE.DIAMOND,
    ]
    colors = ["4F7DF3", "8B5CF6", "F97316", "10B981", "EC4899", "EAB308"]
    for index, (kind, color) in enumerate(zip(kinds, colors, strict=True)):
        shape = _box(
            slide,
            0.05 + index * 2.2,
            1.45,
            2.15,
            1.0,
            kind.name,
            color,
            size=10,
            shape_type=kind,
        )
        shape.rotation = index * 11
        shape.fill.transparency = index * 8
    gradient = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8),
        Inches(4.4),
        Inches(11.7),
        Inches(2.0),
    )
    gradient.fill.gradient()
    gradient.fill.gradient_angle = 35
    gradient.fill.gradient_stops[0].color.rgb = RGBColor.from_string("0D0D2B")
    gradient.fill.gradient_stops[1].color.rgb = RGBColor.from_string("4F7DF3")
    gradient.line.color.rgb = RGBColor.from_string("FFFFFF")
    _text(
        gradient,
        "Linear gradient • transparency • rotation • preset geometry",
        22,
        "FFFFFF",
    )


def _architecture(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title(slide, "System Architecture", "SLIDE 4 / GROUPS + CONNECTORS")
    group = slide.shapes.add_group_shape()
    group.left = Inches(0.7)
    group.top = Inches(1.5)
    group.width = Inches(3.2)
    group.height = Inches(4.8)
    for index, (label, color) in enumerate(
        [("Ingest", "E5EDFF"), ("Parse", "DDF4E8"), ("Resolve", "FFF0D8")]
    ):
        node = group.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.2),
            Inches(0.3 + index * 1.45),
            Inches(2.5),
            Inches(0.9),
        )
        node.fill.solid()
        node.fill.fore_color.rgb = RGBColor.from_string(color)
        _text(node, label, 15)
    group.rotation = 4
    layers = [
        ("Render", 4.8, 2.0, "8B5CF6"),
        ("Browser", 8.0, 2.0, "4F7DF3"),
        ("Evidence", 10.2, 4.7, "10B981"),
    ]
    for label, left, top, color in layers:
        _box(slide, left, top, 2.35, 1.25, label, color, size=17)
    paths = [
        ((3.6, 3.5), (4.8, 2.65), MSO_CONNECTOR.STRAIGHT),
        ((7.15, 2.65), (8.0, 2.65), MSO_CONNECTOR.STRAIGHT),
        ((9.2, 3.25), (10.2, 4.9), MSO_CONNECTOR.CURVE),
    ]
    for start, end, connector_type in paths:
        connector = slide.shapes.add_connector(
            connector_type,
            Inches(start[0]),
            Inches(start[1]),
            Inches(end[0]),
            Inches(end[1]),
        )
        connector.line.color.rgb = RGBColor.from_string("667085")
        connector.line.width = Pt(2.5)


def _table(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title(slide, "Delivery Matrix", "SLIDE 5 / TABLE")
    shape = slide.shapes.add_table(
        6, 5, Inches(0.65), Inches(1.3), Inches(12.0), Inches(5.6)
    )
    table = shape.table
    values = [
        ["Workstream", "Owner", "Status", "Confidence", "Next gate"],
        ["Parser", "Ari", "Ready", "92%", "Corpus"],
        ["Resolver", "Bo", "Active", "81%", "Inheritance"],
        ["Renderer", "Chen", "Active", "76%", "Pixel diff"],
        ["Bindings", "Dia", "Ready", "95%", "Smoke"],
        ["Release", "Eli", "Blocked", "64%", "Native PNG"],
    ]
    for row_index, row in enumerate(values):
        for column_index, value in enumerate(row):
            cell = table.cell(row_index, column_index)
            cell.text = value
            cell.fill.solid()
            fill = "172033" if row_index == 0 else (
                "F3F6FC" if row_index % 2 else "E5EDFF"
            )
            cell.fill.fore_color.rgb = RGBColor.from_string(fill)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.bold = row_index == 0
                paragraph.font.color.rgb = RGBColor.from_string(
                    "FFFFFF" if row_index == 0 else "172033"
                )
                paragraph.alignment = PP_ALIGN.CENTER
    table.cell(4, 0).merge(table.cell(5, 0))


def _images(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title(slide, "Media Treatment", "SLIDE 6 / IMAGE CROP")
    images = [
        ("A", "4F7DF3", "8B5CF6"),
        ("B", "F97316", "FACC15"),
        ("C", "10B981", "22D3EE"),
    ]
    for index, (label, first, second) in enumerate(images):
        picture = slide.shapes.add_picture(
            io.BytesIO(_sample_image(label, first, second)),
            Inches(0.7 + index * 4.2),
            Inches(1.3),
            Inches(3.75),
            Inches(4.7),
        )
        picture.crop_left = 0.08 * index
        picture.crop_right = 0.04 * index
        picture.rotation = (-6, 4, -3)[index]
        picture.line.color.rgb = RGBColor.from_string("FFFFFF")
        picture.line.width = Pt(4)
    _box(
        slide,
        3.7,
        6.2,
        5.9,
        0.75,
        "Embedded PNG • crop • rotation • border",
        "172033",
        size=14,
        text_color="FFFFFF",
    )


def _chart(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title(slide, "Trend & Mix", "SLIDE 7 / CHART")
    data = ChartData()
    data.categories = ["Jan", "Feb", "Mar", "Apr", "May", "Jun"]
    data.add_series("Actual", (18, 25, 21, 34, 42, 48))
    data.add_series("Plan", (20, 23, 27, 31, 38, 44))
    slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(0.7),
        Inches(1.25),
        Inches(8.0),
        Inches(5.8),
        data,
    )
    _box(slide, 9.1, 1.4, 3.4, 1.35, "48\nCurrent", "DDF4E8", size=18)
    _box(slide, 9.1, 3.15, 3.4, 1.35, "+16%\nGrowth", "E5EDFF", size=18)
    _box(slide, 9.1, 4.9, 3.4, 1.35, "2\nExceptions", "FFF0D8", size=18)


def _combined(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = _box(
        slide,
        0,
        0,
        13.333,
        7.5,
        "",
        "0D0D2B",
        shape_type=MSO_SHAPE.RECTANGLE,
    )
    background.line.fill.background()
    _title(slide, "Release Control Room", "SLIDE 8 / COMBINED")
    for shape in [slide.shapes[-2], slide.shapes[-1]]:
        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.color.rgb = RGBColor.from_string("FFFFFF")
    _box(slide, 0.65, 1.2, 3.5, 1.25, "Build\nGREEN", "DDF4E8", size=18)
    _box(slide, 4.5, 1.2, 3.5, 1.25, "Coverage\nCOMPLETE", "E5EDFF", size=18)
    _box(
        slide,
        8.35,
        1.2,
        4.3,
        1.25,
        "Pixel exact\nAWAITING NATIVE",
        "FFF0D8",
        size=16,
    )
    stages = [
        ("Parse", 0.94, "10B981"),
        ("Resolve", 0.86, "4F7DF3"),
        ("Render", 0.72, "8B5CF6"),
    ]
    for index, (label, value, color) in enumerate(stages):
        _box(slide, 0.7, 3.0 + index * 1.05, 2.1, 0.65, label, "172033", size=13)
        track = _box(
            slide,
            3.0,
            3.12 + index * 1.05,
            8.7,
            0.35,
            "",
            "344054",
            shape_type=MSO_SHAPE.RECTANGLE,
        )
        track.line.fill.background()
        progress = _box(
            slide,
            3.0,
            3.12 + index * 1.05,
            8.7 * value,
            0.35,
            "",
            color,
            shape_type=MSO_SHAPE.RECTANGLE,
        )
        progress.line.fill.background()
    _box(slide, 9.6, 6.45, 3.0, 0.55, "STRICT GATE: ARMED", "FFFFFF", size=11)


def create_stress_deck(output_dir: Path) -> list[Path]:
    presentation = Presentation()
    presentation.slide_width = Inches(40 / 3)
    presentation.slide_height = Inches(7.5)
    builders = [
        _dashboard,
        _text_layout,
        _shapes,
        _architecture,
        _table,
        _images,
        _chart,
        _combined,
    ]
    for build_slide in builders:
        build_slide(presentation)
    output = output_dir / "stress_01_complex.pptx"
    presentation.save(output)
    return [output]
