from __future__ import annotations

import argparse
import json
import logging
from collections import Counter
from collections.abc import Sequence
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from evaluate.adjustment_cases import (
    AdjustmentCase,
    build_adjustment_cases,
    load_adjustment_specs,
)
from evaluate.create_adjustment_benchmark_deck import (
    apply_adjustments,
    canonicalize_pptx,
)

logger = logging.getLogger(__name__)

PAIRS_PER_SLIDE = 4
VARIANTS_PER_PAIR = 3
CASES_PER_SLIDE = PAIRS_PER_SLIDE * VARIANTS_PER_PAIR


def write_exhaustive_adjustment_deck(
    adjustment_manifest: Path,
    output_dir: Path,
) -> tuple[Path, Path]:
    specs = load_adjustment_specs(adjustment_manifest)
    cases = build_adjustment_cases(specs)
    output_dir.mkdir(parents=True, exist_ok=True)
    deck_path = output_dir / "all_adjustments.pptx"
    manifest_path = output_dir / "manifest.json"

    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    entries: list[dict[str, object]] = []
    for offset in range(0, len(cases), CASES_PER_SLIDE):
        slide_cases = cases[offset : offset + CASES_PER_SLIDE]
        slide_index = len(presentation.slides)
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        entries.extend(
            _populate_slide(
                slide,
                slide_cases,
                case_offset=offset,
                slide_index=slide_index,
            )
        )

    presentation.save(str(deck_path))
    canonicalize_pptx(deck_path)
    payload = {
        "schema_version": 1,
        "pptx": deck_path.name,
        "source_manifest": adjustment_manifest.name,
        "adjustment_pair_count": len(specs),
        "case_count": len(cases),
        "slide_count": len(presentation.slides),
        "variants": ["low", "default", "high"],
        "pairs_per_slide": PAIRS_PER_SLIDE,
        "range_verification_counts": dict(
            Counter(case.range_verification for case in cases)
        ),
        "entries": entries,
    }
    _ = manifest_path.write_text(
        json.dumps(payload, indent=2, sort_keys=True) + "\n",
        encoding="utf-8",
    )
    return deck_path, manifest_path


def _populate_slide(
    slide: object,
    cases: tuple[AdjustmentCase, ...],
    *,
    case_offset: int,
    slide_index: int,
) -> list[dict[str, object]]:
    entries: list[dict[str, object]] = []
    for local_index, case in enumerate(cases):
        row = local_index // VARIANTS_PER_PAIR
        column = local_index % VARIANTS_PER_PAIR
        row_top = 0.08 + row * 1.85
        cell_left = 2.18 + column * 3.68
        cell_top = row_top + 0.32
        cell_width = 3.48
        cell_height = 1.42
        shape_width = 2.45
        shape_height = 1.08
        shape_left = cell_left + (cell_width - shape_width) / 2
        shape_top = cell_top + (cell_height - shape_height) / 2
        pair_index = (case_offset + local_index) // VARIANTS_PER_PAIR
        shape_name = f"ADJ_{pair_index:03d}_{case.variant.upper()}"

        if column == 0:
            _add_label(
                slide,
                0.08,
                row_top + 0.52,
                2.02,
                0.56,
                f"{case.preset}\n{case.key}",
                size=8,
            )
        _add_label(
            slide,
            cell_left,
            row_top,
            cell_width,
            0.26,
            f"{case.variant.upper()} {case.value}",
            size=7,
        )
        connector_type = _connector_type(case.preset)
        if connector_type is None:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(shape_left),
                Inches(shape_top),
                Inches(shape_width),
                Inches(shape_height),
            )
            shape._element.nvSpPr.cNvPr.set("name", shape_name)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0x44, 0x72, 0xC4)
        else:
            shape = slide.shapes.add_connector(
                connector_type,
                Inches(shape_left),
                Inches(shape_top),
                Inches(shape_left + shape_width),
                Inches(shape_top + shape_height),
            )
            shape._element.nvCxnSpPr.cNvPr.set("name", shape_name)
        shape._element.spPr.prstGeom.set("prst", case.preset)
        shape.line.color.rgb = RGBColor(0x20, 0x20, 0x20)
        shape.line.width = Pt(1.25)
        apply_adjustments(shape, case.adjustments)
        entries.append(
            {
                "shape_name": shape_name,
                "pair_index": pair_index,
                "preset": case.preset,
                "key": case.key,
                "variant": case.variant,
                "value": case.value,
                "adjustments": case.adjustments,
                "range_verification": case.range_verification,
                "slide_index": slide_index,
                "slot_index": local_index,
                "crop_in": {
                    "left": round(cell_left, 4),
                    "top": round(cell_top, 4),
                    "width": cell_width,
                    "height": cell_height,
                },
                "shape_box_in": {
                    "left": round(shape_left, 4),
                    "top": round(shape_top, 4),
                    "width": shape_width,
                    "height": shape_height,
                },
            }
        )
    return entries


def _connector_type(preset: str):
    if preset.startswith("bentConnector"):
        return MSO_CONNECTOR.ELBOW
    if preset.startswith("curvedConnector"):
        return MSO_CONNECTOR.CURVE
    if preset == "straightConnector1":
        return MSO_CONNECTOR.STRAIGHT
    return None


def _add_label(
    slide: object,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    size: int,
) -> None:
    shape = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = RGBColor(0x17, 0x20, 0x33)


def main(argv: Sequence[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        description="Generate every official preset adjustment variant."
    )
    _ = parser.add_argument(
        "--adjustment-manifest",
        type=Path,
        required=True,
    )
    _ = parser.add_argument("--output-dir", type=Path, required=True)
    args = parser.parse_args(argv)
    deck, manifest = write_exhaustive_adjustment_deck(
        args.adjustment_manifest,
        args.output_dir,
    )
    logger.info("Adjustment deck: %s", deck)
    logger.info("Adjustment manifest: %s", manifest)
    return 0


if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")
    raise SystemExit(main())
