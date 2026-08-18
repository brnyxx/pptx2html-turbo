from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Inches, Pt

if __package__:
    from evaluate.stress_deck import (
        _architecture,
        _box,
        _chart,
        _combined,
        _dashboard,
        _images,
        _shapes,
        _table,
        _text_layout,
        _title,
    )
else:
    from stress_deck import (
        _architecture,
        _box,
        _chart,
        _combined,
        _dashboard,
        _images,
        _shapes,
        _table,
        _text_layout,
        _title,
    )


def _timeline(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _title(slide, "Program Timeline", "SLIDE 9 / CONNECTED FLOW")
    phases = [
        ("Discover", "4F7DF3"),
        ("Design", "8B5CF6"),
        ("Build", "F97316"),
        ("Validate", "10B981"),
        ("Release", "22D3EE"),
    ]
    for index, (label, color) in enumerate(phases):
        left = 0.65 + index * 2.55
        milestone = _box(
            slide,
            left,
            2.15 + (index % 2) * 1.25,
            2.0,
            1.2,
            f"{index + 1:02d}\n{label}",
            color,
            size=15,
        )
        milestone.rotation = -3 if index % 2 else 3
        if index == 0:
            continue
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(left - 0.55),
            Inches(3.05),
            Inches(left),
            Inches(2.75 + (index % 2) * 1.25),
        )
        connector.line.color.rgb = RGBColor.from_string("667085")
        connector.line.width = Pt(2)
    _box(
        slide,
        1.35,
        5.65,
        10.65,
        0.85,
        "Critical path • dependencies • alternating alignment • rotated milestones",
        "F8FAFC",
        size=14,
    )


def _matrix(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _title(slide, "Risk Matrix", "SLIDE 10 / DENSE COMPOSITION")
    colors = ["DDF4E8", "E5EDFF", "FFF0D8", "FFE4E6"]
    for row in range(4):
        for column in range(5):
            score = (row + 1) * (column + 1)
            card = _box(
                slide,
                0.75 + column * 2.45,
                1.25 + row * 1.35,
                2.05,
                1.05,
                f"R{row + 1}.{column + 1}\nScore {score}",
                colors[(row + column) % len(colors)],
                size=12,
            )
            card.rotation = ((row * 5 + column) % 3 - 1) * 1.5
    _box(slide, 9.9, 6.75, 2.45, 0.45, "20 CELLS / 4 LEVELS", "172033", size=9)


def _apply_deck_variant(
    slide: object,
    profile: str,
    deck_index: int,
    slide_index: int,
) -> None:
    title = next(
        shape
        for shape in slide.shapes
        if shape.has_text_frame and shape.text.strip()
    )
    paragraph = title.text_frame.paragraphs[0]
    paragraph.runs[0].text = f"{paragraph.runs[0].text} - {profile}"

    target = slide.shapes[2]
    horizontal_step = deck_index - 5
    vertical_step = ((deck_index * 3 + slide_index) % 10) - 5
    target.left += Inches(horizontal_step * 0.025)
    target.top += Inches(vertical_step * 0.012)
    target.width += Inches(deck_index * 0.006)


def create_challenge_corpus(output_dir: Path) -> list[Path]:
    builders = [
        _dashboard,
        _text_layout,
        _shapes,
        _architecture,
        _table,
        _images,
        _chart,
        _combined,
        _timeline,
        _matrix,
    ]
    accents = [
        ("Atlas", "4F7DF3"),
        ("Beacon", "8B5CF6"),
        ("Cobalt", "F97316"),
        ("Delta", "10B981"),
        ("Ember", "22D3EE"),
        ("Flux", "EC4899"),
        ("Grove", "EAB308"),
        ("Harbor", "14B8A6"),
        ("Ion", "6366F1"),
        ("Junction", "EF4444"),
    ]
    outputs: list[Path] = []
    for deck_index, (profile, accent) in enumerate(accents, start=1):
        presentation = Presentation()
        presentation.slide_width = Inches(40 / 3)
        presentation.slide_height = Inches(7.5)
        for slide_index, build_slide in enumerate(builders, start=1):
            build_slide(presentation)
            slide = presentation.slides[-1]
            _apply_deck_variant(
                slide,
                profile,
                deck_index,
                slide_index,
            )
            badge = _box(
                slide,
                11.85,
                0.08,
                1.05,
                0.3,
                f"C{deck_index:02d}.{slide_index:02d}",
                accent,
                size=7,
            )
            badge.line.fill.background()
            if (deck_index + slide_index) % 2:
                badge.rotation = 2
        output = output_dir / f"challenge_{deck_index:02d}.pptx"
        presentation.save(output)
        outputs.append(output)
    return outputs
