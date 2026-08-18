import contextlib
import json
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from evaluate.completion_deck_package import png_bytes
from evaluate.scaffold_powerpoint_golden_batch import scaffold_powerpoint_golden_batch
from evaluate.summarize_powerpoint_golden import summarize_powerpoint_golden_batch


class SummarizePowerPointGoldenBatchTests(unittest.TestCase):
    def test_reports_complete_batch_and_manifest(self) -> None:
        with self._batch() as (golden, output):
            summary = summarize_powerpoint_golden_batch(golden, output)
        self.assertTrue(summary["evidence_ready_for_exact_promotion"])
        self.assertEqual(summary["provenance_errors"], [])
        self.assertEqual(summary["batch_identity"]["producer"], "Microsoft PowerPoint")

    def test_reports_missing_metadata_and_missing_deck_output(self) -> None:
        with self._batch() as (golden, output):
            (output / "sample/metadata.json").unlink()
            summary = summarize_powerpoint_golden_batch(golden, output)
        self.assertEqual(summary["missing_metadata"], ["sample"])
        self.assertFalse(summary["evidence_ready_for_exact_promotion"])

    def test_reports_manifest_count_mismatch(self) -> None:
        with self._batch() as (golden, output):
            path = output / "manifest.json"
            manifest = json.loads(path.read_text())
            manifest["deck_count"] = 99
            path.write_text(json.dumps(manifest))
            summary = summarize_powerpoint_golden_batch(golden, output)
        self.assertFalse(summary["manifest_deck_count_matches"])
        self.assertFalse(summary["evidence_ready_for_exact_promotion"])

    def test_reports_invalid_metadata_fields_and_bad_slide_names(self) -> None:
        with self._batch() as (golden, output):
            metadata_path = output / "sample/metadata.json"
            metadata = json.loads(metadata_path.read_text())
            metadata["producer"] = "browser"
            metadata_path.write_text(json.dumps(metadata))
            (output / "sample/Slide2.PNG").write_bytes(png_bytes())
            summary = summarize_powerpoint_golden_batch(golden, output)
        self.assertEqual(summary["invalid_metadata"], ["sample"])
        self.assertEqual(summary["incomplete_slide_exports"], ["sample"])

    @contextlib.contextmanager
    def _batch(self):
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            golden, output = root / "golden", root / "output"
            golden.mkdir()
            output.mkdir()
            presentation = Presentation()
            presentation.slides.add_slide(presentation.slide_layouts[6])
            presentation.save(golden / "sample.pptx")
            deck = output / "sample"
            deck.mkdir()
            (deck / "Slide1.PNG").write_bytes(png_bytes())
            scaffold_powerpoint_golden_batch(golden, output, self._metadata())
            yield golden, output

    def _metadata(self):
        return {"producer":"Microsoft PowerPoint","platform":"Windows","powerpoint_version":"16.0.17726.20160","powerpoint_build":"16.0.17726.20160","powerpoint_channel":"Current Channel","windows_version":"Windows 11 23H2","export_command":"pwsh -File reference_render_powerpoint.ps1","output_resolution":"1x1","golden_set_revision":"abc1234","capture_timestamp":"2026-04-02T12:00:00Z","batch_id":"powerpoint-test-batch"}


if __name__ == "__main__":
    unittest.main()
