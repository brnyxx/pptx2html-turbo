import importlib.util
import tempfile
import unittest
from pathlib import Path
from xml.etree import ElementTree
from zipfile import ZipFile

from pptx import Presentation

from evaluate import synthetic_pptx
from evaluate.synthetic_scene import create_synthetic_corpus


class SyntheticPptxTests(unittest.TestCase):
    def test_pptx_emitter_module_exists(self) -> None:
        module = importlib.util.find_spec("evaluate.synthetic_pptx")

        self.assertIsNotNone(module)

    def test_emitter_writes_ten_decks_with_matching_scene_geometry(self) -> None:
        # Given
        corpus = create_synthetic_corpus()
        write_corpus = getattr(
            synthetic_pptx,
            "write_synthetic_pptx_corpus",
            None,
        )

        # When/Then
        self.assertTrue(callable(write_corpus))
        with tempfile.TemporaryDirectory() as tmp:
            paths = write_corpus(corpus, Path(tmp))

            self.assertEqual(
                [path.name for path in paths],
                [f"synthetic_{index:02d}.pptx" for index in range(1, 11)],
            )
            for deck, path in zip(corpus, paths, strict=True):
                presentation = Presentation(path)
                self.assertEqual(len(presentation.slides), 10)
                self.assertEqual(
                    presentation.slide_width,
                    synthetic_pptx.CANVAS_WIDTH_EMU,
                )
                self.assertEqual(
                    presentation.slide_height,
                    synthetic_pptx.CANVAS_HEIGHT_EMU,
                )
                for scene, slide in zip(
                    deck.scenes,
                    presentation.slides,
                    strict=True,
                ):
                    self.assertEqual(len(slide.shapes), len(scene.rectangles))
                    for rectangle, shape in zip(
                        scene.rectangles,
                        slide.shapes,
                        strict=True,
                    ):
                        self.assertEqual(
                            (
                                shape.left,
                                shape.top,
                                shape.width,
                                shape.height,
                            ),
                            (
                                rectangle.x,
                                rectangle.y,
                                rectangle.width,
                                rectangle.height,
                            ),
                        )
                        self.assertEqual(
                            str(shape.fill.fore_color.rgb),
                            rectangle.fill,
                        )
                        self.assertEqual(shape.text, "")

    def test_cli_writes_corpus_to_requested_directory(self) -> None:
        # Given
        main = getattr(synthetic_pptx, "main", None)

        # When/Then
        self.assertTrue(callable(main))
        with tempfile.TemporaryDirectory() as tmp:
            output = Path(tmp) / "decks"
            exit_code = main(["--output", str(output)])

            self.assertEqual(exit_code, 0)
            self.assertEqual(len(list(output.glob("synthetic_*.pptx"))), 10)

    def test_emitter_disables_theme_effect_reference(self) -> None:
        # Given
        write_corpus = synthetic_pptx.write_synthetic_pptx_corpus
        corpus = create_synthetic_corpus()[:1]

        # When
        with tempfile.TemporaryDirectory() as tmp:
            path = write_corpus(corpus, Path(tmp))[0]
            with ZipFile(path) as archive:
                slide_xml = ElementTree.fromstring(
                    archive.read("ppt/slides/slide1.xml")
                )
                self.assertIsNone(archive.testzip())
            self.assertFalse(path.with_suffix(".pptx.tmp").exists())

        # Then
        effect_refs = slide_xml.findall(
            ".//a:effectRef",
            {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"},
        )
        self.assertEqual(
            [effect.attrib["idx"] for effect in effect_refs],
            ["0"] * 8,
        )

    def test_effect_normalizer_preserves_source_when_marker_is_missing(self) -> None:
        # Given
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "missing-effect.pptx"
            with ZipFile(path, "w") as archive:
                archive.writestr(
                    "ppt/slides/slide1.xml",
                    b"<p:sld xmlns:p='urn:p'><p:cSld/></p:sld>",
                )
            original = path.read_bytes()

            # When/Then
            with self.assertRaisesRegex(
                ValueError,
                "SYNTHETIC_EFFECT_REFERENCE_MISSING",
            ):
                synthetic_pptx._disable_theme_effects(path)
            self.assertEqual(path.read_bytes(), original)
            self.assertFalse(path.with_suffix(".pptx.tmp").exists())


if __name__ == "__main__":
    unittest.main()
