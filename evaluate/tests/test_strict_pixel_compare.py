import contextlib
import io
import json
import tempfile
import unittest
from pathlib import Path

from PIL import Image
from pptx import Presentation

from evaluate.scaffold_powerpoint_golden_batch import (
    scaffold_powerpoint_golden_batch,
)
from evaluate.strict_pixel_compare import (
    StrictPixelError,
    compare_strict_batch,
    main,
)


class StrictPixelCompareTests(unittest.TestCase):
    def test_identical_powerpoint_and_browser_pixels_are_exact(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            golden_set, references, candidates = self._fixture(Path(tmp))

            result = compare_strict_batch(golden_set, references, candidates)

            self.assertTrue(result["ok"])
            self.assertEqual(result["deck_count"], 1)
            self.assertEqual(result["slide_count"], 1)
            self.assertEqual(result["total_pixels"], 2)
            self.assertEqual(result["mismatched_pixels"], 0)
            self.assertEqual(result["max_channel_delta"], 0)
            self.assertEqual(result["slides"][0]["size"], [2, 1])
            self.assertTrue(result["slides"][0]["exact"])

    def test_one_channel_delta_in_one_pixel_fails_exact_gate(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            golden_set, references, candidates = self._fixture(
                root,
                candidate_pixels=((255, 0, 0, 255), (0, 254, 0, 255)),
            )
            diff_dir = root / "diffs"

            result = compare_strict_batch(
                golden_set,
                references,
                candidates,
                diff_dir=diff_dir,
            )

            self.assertFalse(result["ok"])
            self.assertEqual(result["mismatched_pixels"], 1)
            self.assertEqual(result["max_channel_delta"], 1)
            self.assertEqual(result["slides"][0]["mismatch_ratio"], 0.5)
            self.assertTrue((diff_dir / "sample" / "Slide1.PNG").is_file())

    def test_one_blue_channel_step_is_not_lost_to_grayscale_rounding(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            golden_set, references, candidates = self._fixture(
                Path(tmp),
                candidate_pixels=((255, 0, 0, 255), (0, 255, 1, 255)),
            )

            result = compare_strict_batch(golden_set, references, candidates)

            self.assertFalse(result["ok"])
            self.assertEqual(result["mismatched_pixels"], 1)
            self.assertEqual(result["max_channel_delta"], 1)

    def test_dimension_mismatch_is_rejected_without_resizing(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            golden_set, references, candidates = self._fixture(root)
            self._save_png(
                candidates / "sample" / "slide_0.png",
                (1, 1),
                ((255, 0, 0, 255),),
            )

            with self.assertRaisesRegex(
                StrictPixelError,
                "PIXEL_DIMENSION_MISMATCH",
            ):
                compare_strict_batch(golden_set, references, candidates)

    def test_missing_candidate_slide_is_rejected(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            golden_set, references, candidates = self._fixture(Path(tmp))
            (candidates / "sample" / "slide_0.png").unlink()

            with self.assertRaisesRegex(
                StrictPixelError,
                "PIXEL_CANDIDATE_MISSING",
            ):
                compare_strict_batch(golden_set, references, candidates)

    def test_cli_returns_one_and_emits_json_for_pixel_difference(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            golden_set, references, candidates = self._fixture(
                root,
                candidate_pixels=((255, 0, 0, 255), (0, 254, 0, 255)),
            )
            output = root / "result.json"

            with contextlib.redirect_stdout(io.StringIO()):
                exit_code = main(
                    [
                        "--golden-set-dir",
                        str(golden_set),
                        "--reference-dir",
                        str(references),
                        "--candidate-dir",
                        str(candidates),
                        "--output-json",
                        str(output),
                    ]
                )

            self.assertEqual(exit_code, 1)
            self.assertFalse(json.loads(output.read_text(encoding="utf-8"))["ok"])

    def _fixture(
        self,
        root: Path,
        *,
        candidate_pixels: tuple[tuple[int, int, int, int], ...] = (
            (255, 0, 0, 255),
            (0, 255, 0, 255),
        ),
    ) -> tuple[Path, Path, Path]:
        golden_set = root / "golden_set"
        references = root / "powerpoint_golden"
        candidates = root / "candidates"
        golden_set.mkdir()
        references.mkdir()
        (references / "sample").mkdir()
        (candidates / "sample").mkdir(parents=True)

        presentation = Presentation()
        presentation.slides.add_slide(presentation.slide_layouts[6])
        presentation.save(golden_set / "sample.pptx")
        pixels = ((255, 0, 0, 255), (0, 255, 0, 255))
        self._save_png(references / "sample" / "Slide1.PNG", (2, 1), pixels)
        self._save_png(
            candidates / "sample" / "slide_0.png",
            (2, 1),
            candidate_pixels,
        )
        scaffold_powerpoint_golden_batch(
            golden_set,
            references,
            {
                "producer": "Microsoft PowerPoint",
                "platform": "Windows",
                "powerpoint_version": "16.0",
                "powerpoint_build": "16.0.17726.20160",
                "powerpoint_channel": "Current Channel",
                "windows_version": "Windows 11 23H2",
                "export_command": "pwsh -File reference_render_powerpoint.ps1",
                "output_resolution": "2x1",
                "golden_set_revision": "abc1234",
                "capture_timestamp": "2026-08-13T12:00:00Z",
                "batch_id": "powerpoint-strict-pixel-test",
            },
        )
        return golden_set, references, candidates

    @staticmethod
    def _save_png(
        path: Path,
        size: tuple[int, int],
        pixels: tuple[tuple[int, int, int, int], ...],
    ) -> None:
        image = Image.new("RGBA", size)
        image.putdata(pixels)
        image.save(path)


if __name__ == "__main__":
    unittest.main()
