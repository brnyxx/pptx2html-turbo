import json
import re
import tempfile
import unittest
from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation

from evaluate.challenge_deck import create_challenge_corpus

MANIFEST = Path(__file__).parents[1] / "visual_element_coverage.json"
SHAPE_MODEL = (
    Path(__file__).parents[2]
    / "crates"
    / "pptx2html-core"
    / "src"
    / "model"
    / "shape.rs"
)
REQUIRED_USER_ELEMENTS = {
    "arrow",
    "diamond",
    "curved-connector",
    "straight-connector",
    "chart",
    "table",
}
ALLOWED_TIERS = {
    "synthetic-exact",
    "challenge-proxy",
    "rust-regression",
    "fallback-contract",
}
CHALLENGE_MARKERS = {
    "arrow": 'prst="rightArrow"',
    "border": "<a:ln",
    "chart": "<c:lineChart>",
    "curved-connector": 'prst="curvedConnector3"',
    "custom-preset": 'prst="chevron"',
    "diamond": 'prst="diamond"',
    "gradient-fill": "<a:gradFill",
    "group": "<p:grpSp>",
    "image-crop": "<a:srcRect",
    "rotation-flip": ' rot="',
    "rounded-rectangle": 'prst="roundRect"',
    "straight-connector": 'prst="line"',
    "table": "<a:tbl>",
    "text-typography": "<a:t>",
}
SYNTHETIC_EXACT_ELEMENTS = {
    "overlap-order",
    "slide-background",
    "solid-fill",
    "solid-rectangle",
}


class VisualElementCoverageTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls) -> None:
        cls.temporary = tempfile.TemporaryDirectory()
        output = Path(cls.temporary.name)
        deck = create_challenge_corpus(output)[0]
        cls.deck = deck
        with ZipFile(deck) as archive:
            cls.parts = {
                name: archive.read(name).decode(errors="ignore")
                for name in archive.namelist()
                if name.endswith(".xml")
            }
        cls.challenge_payload = "\n".join(cls.parts.values())

    @classmethod
    def tearDownClass(cls) -> None:
        cls.temporary.cleanup()

    def test_visual_element_manifest_exists(self) -> None:
        self.assertTrue(MANIFEST.is_file())

    def test_manifest_covers_every_shape_type_variant(self) -> None:
        # Given
        source = SHAPE_MODEL.read_text(encoding="utf-8")
        enum_body = source.split("pub enum ShapeType {", maxsplit=1)[1].split(
            "\n}",
            maxsplit=1,
        )[0]
        variants = set(
            re.findall(r"^\s+([A-Z][A-Za-z0-9]+)(?:\(|,)", enum_body, re.MULTILINE)
        )

        # When
        manifest = json.loads(MANIFEST.read_text(encoding="utf-8"))

        # Then
        self.assertEqual(set(manifest["shape_type_variants"]), variants)

    def test_manifest_has_evidence_for_every_required_visual_element(self) -> None:
        # Given/When
        manifest = json.loads(MANIFEST.read_text(encoding="utf-8"))
        elements = manifest["elements"]
        element_ids = [element["id"] for element in elements]

        # Then
        self.assertEqual(len(element_ids), len(set(element_ids)))
        self.assertTrue(REQUIRED_USER_ELEMENTS.issubset(element_ids))
        for element in elements:
            self.assertTrue(element["evidence"], element["id"])
            for evidence in element["evidence"]:
                self.assertIn(evidence["tier"], ALLOWED_TIERS)
                self.assertTrue(evidence["source"])

    def test_manifest_shape_mappings_reference_declared_elements(self) -> None:
        manifest = json.loads(MANIFEST.read_text(encoding="utf-8"))
        element_ids = {element["id"] for element in manifest["elements"]}

        self.assertTrue(
            set(manifest["shape_type_variants"].values()).issubset(element_ids)
        )

    def test_every_challenge_evidence_has_an_ooxml_probe(self) -> None:
        manifest = json.loads(MANIFEST.read_text(encoding="utf-8"))
        challenge_ids = {
            element["id"]
            for element in manifest["elements"]
            if any(
                evidence["tier"] == "challenge-proxy"
                for evidence in element["evidence"]
            )
        }

        self.assertEqual(challenge_ids, set(CHALLENGE_MARKERS))
        for element_id, marker in CHALLENGE_MARKERS.items():
            self.assertIn(marker, self.challenge_payload, element_id)

    def test_every_synthetic_exact_evidence_maps_to_scene_contract(self) -> None:
        manifest = json.loads(MANIFEST.read_text(encoding="utf-8"))
        exact_ids = {
            element["id"]
            for element in manifest["elements"]
            if any(
                evidence["tier"] == "synthetic-exact"
                for evidence in element["evidence"]
            )
        }

        self.assertEqual(exact_ids, SYNTHETIC_EXACT_ELEMENTS)

    def test_every_rust_regression_evidence_identifier_exists(self) -> None:
        manifest = json.loads(MANIFEST.read_text(encoding="utf-8"))
        rust_root = SHAPE_MODEL.parents[2]
        test_functions = {
            function
            for path in rust_root.rglob("*.rs")
            for function in re.findall(
                r"#\[test\]\s+fn\s+([a-zA-Z0-9_]+)",
                path.read_text(encoding="utf-8"),
            )
        }
        identifiers = {
            evidence["source"].rsplit(":", maxsplit=1)[-1]
            for element in manifest["elements"]
            for evidence in element["evidence"]
            if evidence["tier"] == "rust-regression"
        }

        for identifier in identifiers:
            self.assertIn(identifier, test_functions)

    def test_new_shape_evidence_links_to_stress_deck_source(self) -> None:
        manifest = json.loads(MANIFEST.read_text(encoding="utf-8"))
        elements = {element["id"]: element for element in manifest["elements"]}

        for element_id in ("arrow", "curved-connector", "diamond"):
            sources = [
                evidence["source"]
                for evidence in elements[element_id]["evidence"]
                if evidence["tier"] == "challenge-proxy"
            ]
            self.assertTrue(
                any("stress_deck" in source for source in sources),
                element_id,
            )

    def test_challenge_deck_contains_arrow_element(self) -> None:
        self.assertIn('prst="rightArrow"', self.parts["ppt/slides/slide3.xml"])

    def test_challenge_deck_contains_diamond_element(self) -> None:
        self.assertIn('prst="diamond"', self.parts["ppt/slides/slide3.xml"])

    def test_challenge_deck_contains_curved_connector(self) -> None:
        self.assertIn(
            'prst="curvedConnector3"',
            self.parts["ppt/slides/slide4.xml"],
        )

    def test_challenge_deck_contains_straight_connector(self) -> None:
        self.assertIn(
            'prst="line"',
            self.parts["ppt/slides/slide4.xml"],
        )

    def test_challenge_deck_contains_chart(self) -> None:
        chart_parts = [
            payload
            for name, payload in self.parts.items()
            if name.startswith("ppt/charts/chart")
        ]
        self.assertTrue(any("<c:lineChart>" in payload for payload in chart_parts))

    def test_challenge_deck_contains_table(self) -> None:
        self.assertIn("<a:tbl>", self.parts["ppt/slides/slide5.xml"])

    def test_challenge_deck_contains_text_typography(self) -> None:
        self.assertIn("<a:t>", self.parts["ppt/slides/slide2.xml"])

    def test_challenge_deck_contains_image_crop(self) -> None:
        self.assertIn("<a:srcRect", self.parts["ppt/slides/slide6.xml"])

    def test_challenge_deck_contains_gradient_fill(self) -> None:
        self.assertIn("<a:gradFill", self.parts["ppt/slides/slide3.xml"])

    def test_challenge_deck_contains_group_shape(self) -> None:
        self.assertIn("<p:grpSp>", self.parts["ppt/slides/slide4.xml"])

    def test_challenge_deck_contains_preset_shape_family(self) -> None:
        slide = self.parts["ppt/slides/slide3.xml"]
        for preset in (
            "arc",
            "chevron",
            "diamond",
            "hexagon",
            "rightArrow",
            "star5",
        ):
            self.assertIn(f'prst="{preset}"', slide)

    def test_chevron_has_non_degenerate_authored_text_area(self) -> None:
        presentation = Presentation(str(self.deck))
        shape_slide = presentation.slides[2]
        chevron = next(
            shape
            for shape in shape_slide.shapes
            if shape.has_text_frame and shape.text == "CHEVRON"
        )

        self.assertGreaterEqual(chevron.width / chevron.height, 2.0)

    def test_text_stress_uses_explicit_visible_text_colors(self) -> None:
        text_slide = self.parts["ppt/slides/slide2.xml"]
        media_slide = self.parts["ppt/slides/slide6.xml"]
        caption_index = media_slide.index("<a:t>Embedded PNG")
        caption_markup = media_slide[caption_index - 300 : caption_index + 100]

        self.assertGreaterEqual(text_slide.count('val="172033"'), 17)
        self.assertIn('val="FFFFFF"', caption_markup)


if __name__ == "__main__":
    unittest.main()
