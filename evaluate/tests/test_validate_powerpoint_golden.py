import contextlib
import binascii
import json
import struct
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from evaluate.completion_deck_package import png_bytes
from evaluate.powerpoint_provenance import (
    ProvenanceError,
    validate_png,
    validate_provenance,
)
from evaluate.scaffold_powerpoint_golden_batch import scaffold_powerpoint_golden_batch
from evaluate.validate_powerpoint_golden import ValidationError, validate_powerpoint_golden_batch


class ValidatePowerPointGoldenBatchTests(unittest.TestCase):
    def test_rejects_non_powerpoint_export_command(self) -> None:
        metadata = self._metadata()
        metadata["export_command"] = "libreoffice --headless --convert-to png"
        self.assertIn("PROVENANCE_EXPORT_COMMAND_INVALID", validate_provenance(metadata))

    def test_rejects_crc_correct_png_with_invalid_idat(self) -> None:
        def chunk(kind: bytes, payload: bytes) -> bytes:
            checksum = binascii.crc32(kind + payload) & 0xFFFFFFFF
            return struct.pack(">I", len(payload)) + kind + payload + struct.pack(">I", checksum)

        data = (
            b"\x89PNG\r\n\x1a\n"
            + chunk(b"IHDR", struct.pack(">IIBBBBB", 1, 1, 8, 6, 0, 0, 0))
            + chunk(b"IDAT", b"not-zlib")
            + chunk(b"IEND", b"")
        )
        with tempfile.TemporaryDirectory() as temporary:
            path = Path(temporary) / "invalid.png"
            path.write_bytes(data)
            with self.assertRaisesRegex(ProvenanceError, "PNG_IDAT_INVALID"):
                validate_png(path)

    def test_accepts_complete_batch_and_returns_summary(self) -> None:
        with self._batch() as (golden, output):
            summary = validate_powerpoint_golden_batch(golden, output)
        self.assertEqual(summary["deck_count"], 1)
        self.assertEqual(summary["slide_image_count"], 1)

    def test_rejects_missing_metadata_manifest(self) -> None:
        with self._batch() as (golden, output):
            (output / "manifest.json").unlink()
            with self.assertRaisesRegex(ValidationError, "JSON_FILE_INVALID"):
                validate_powerpoint_golden_batch(golden, output)

    def test_rejects_missing_slide_export(self) -> None:
        with self._batch() as (golden, output):
            (output / "sample/Slide1.PNG").unlink()
            with self.assertRaisesRegex(ValidationError, "PNG_FILE_INVALID"):
                validate_powerpoint_golden_batch(golden, output)

    def test_rejects_png_crc_and_cross_link_hash_mutations(self) -> None:
        with self._batch() as (golden, output):
            image = output / "sample/Slide1.PNG"
            data = bytearray(image.read_bytes())
            data[-5] ^= 1
            image.write_bytes(data)
            with self.assertRaisesRegex(ValidationError, "PNG_CRC_INVALID"):
                validate_powerpoint_golden_batch(golden, output)
        with self._batch() as (golden, output):
            metadata_path = output / "sample/metadata.json"
            metadata = json.loads(metadata_path.read_text())
            metadata["source_sha256"] = "0" * 64
            metadata_path.write_text(json.dumps(metadata))
            with self.assertRaisesRegex(ValidationError, "SOURCE_SHA256_MISMATCH"):
                validate_powerpoint_golden_batch(golden, output)

    def test_reviewer_fake_native_mutation_is_rejected(self) -> None:
        """PowerPoint-native oracle text plus LibreOffice/browser bytes is not native evidence."""
        with self._batch() as (golden, output):
            metadata_path = output / "sample/metadata.json"
            metadata = json.loads(metadata_path.read_text())
            metadata.update(
                {
                    "oracle": "PowerPoint-native",
                    "producer": "LibreOffice",
                    "platform": "browser",
                    "powerpoint_version": "16.0.0",
                    "powerpoint_build": "16.0.0",
                }
            )
            metadata_path.write_text(json.dumps(metadata))
            (output / "sample/Slide1.PNG").write_bytes(b"fabricated secondary bytes")
            with self.assertRaisesRegex(
                ValidationError,
                "PROVENANCE_PRODUCER_NOT_POWERPOINT|PNG_SIGNATURE_INVALID",
            ):
                validate_powerpoint_golden_batch(golden, output)

    @contextlib.contextmanager
    def _batch(self):
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            golden, output = root / "golden", root / "output"
            golden.mkdir()
            output.mkdir()
            presentation = Presentation()
            presentation.slides.add_slide(presentation.slide_layouts[6])
            presentation.save(golden / "sample.pptx")
            deck = output / "sample"
            deck.mkdir()
            (deck / "Slide1.PNG").write_bytes(png_bytes())
            scaffold_powerpoint_golden_batch(golden, output, self._metadata())
            yield golden, output

    def _metadata(self) -> dict[str, str]:
        return {
            "producer": "Microsoft PowerPoint",
            "platform": "Windows",
            "powerpoint_version": "16.0.17726.20160",
            "powerpoint_build": "16.0.17726.20160",
            "powerpoint_channel": "Current Channel",
            "windows_version": "Windows 11 23H2",
            "export_command": "pwsh -File reference_render_powerpoint.ps1",
            "output_resolution": "960x540",
            "golden_set_revision": "abc1234",
            "capture_timestamp": "2026-04-02T12:00:00Z",
            "batch_id": "powerpoint-test-batch",
        }


if __name__ == "__main__":
    unittest.main()
