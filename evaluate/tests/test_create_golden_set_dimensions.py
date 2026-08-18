import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

from pptx import Presentation

from evaluate import create_golden_set as golden_generator
from evaluate.create_golden_set import create_golden_set


def _visible_slide_signatures(
    presentation: Presentation,
) -> list[tuple[tuple[str, ...], tuple[tuple[int, ...], ...]]]:
    signatures = []
    for slide in presentation.slides:
        texts: list[str] = []
        geometry: list[tuple[int, ...]] = []
        for shape in slide.shapes:
            text = shape.text.strip() if shape.has_text_frame else ""
            if (
                len(text) == 6
                and text.startswith("C")
                and text[1:3].isdigit()
                and text[3] == "."
                and text[4:6].isdigit()
            ):
                continue
            if shape.has_table:
                text = "\n".join(
                    cell.text
                    for row in shape.table.rows
                    for cell in row.cells
                )
            texts.append(text)
            geometry.append(
                (
                    int(shape.shape_type),
                    shape.left,
                    shape.top,
                    shape.width,
                    shape.height,
                    round(shape.rotation * 100),
                )
            )
        signatures.append((tuple(texts), tuple(geometry)))
    return signatures


class CreateGoldenSetDimensionsTests(unittest.TestCase):
    def test_full_corpus_uses_960_by_540_aspect_ratio(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            output = Path(tmp)

            generated_by_category = create_golden_set(
                output,
                categories=["basic_text"],
            )
            generated = generated_by_category["basic_text"]

            self.assertTrue(generated)
            for path in generated:
                presentation = Presentation(path)
                self.assertEqual(
                    presentation.slide_width * 540,
                    presentation.slide_height * 960,
                    path.name,
                )

    def test_generation_failure_is_not_silently_accepted(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            output = Path(tmp)

            with (
                patch.dict(
                golden_generator.CATEGORY_GENERATORS,
                {"broken": lambda _: (_ for _ in ()).throw(RuntimeError("broken"))},
                clear=True,
                ),
                self.assertRaisesRegex(RuntimeError, "broken"),
            ):
                create_golden_set(output)

    def test_challenge_category_generates_ten_decks_with_ten_slides_each(
        self,
    ) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            output = Path(tmp)

            generated_by_category = create_golden_set(
                output,
                categories=["challenge"],
            )
            generated = generated_by_category["challenge"]

            self.assertEqual(len(generated), 10)
            self.assertEqual(
                [path.name for path in generated],
                [f"challenge_{index:02d}.pptx" for index in range(1, 11)],
            )
            for path in generated:
                presentation = Presentation(path)
                self.assertEqual(len(presentation.slides), 10, path.name)
                self.assertEqual(
                    presentation.slide_width * 540,
                    presentation.slide_height * 960,
                    path.name,
                )
                self.assertGreater(
                    sum(len(slide.shapes) for slide in presentation.slides),
                    60,
                    path.name,
                )

    def test_challenge_slides_are_not_badge_only_variants(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            generated = create_golden_set(
                Path(tmp),
                categories=["challenge"],
            )["challenge"]

            signatures = [
                _visible_slide_signatures(Presentation(path))
                for path in generated
            ]
            flattened = [
                signature
                for deck_signatures in signatures
                for signature in deck_signatures
            ]

            self.assertEqual(len(set(flattened)), 100)
            for slide_index in range(10):
                texts = {
                    deck_signatures[slide_index][0]
                    for deck_signatures in signatures
                }
                geometry = {
                    deck_signatures[slide_index][1]
                    for deck_signatures in signatures
                }
                self.assertEqual(len(texts), 10, slide_index)
                self.assertEqual(len(geometry), 10, slide_index)

    def test_challenge_profiles_only_extend_nonempty_titles(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            generated = create_golden_set(
                Path(tmp),
                categories=["challenge"],
            )["challenge"]

            for path in generated:
                presentation = Presentation(path)
                for slide_index, slide in enumerate(presentation.slides):
                    texts = [
                        shape.text.strip()
                        for shape in slide.shapes
                        if shape.has_text_frame
                    ]
                    self.assertFalse(
                        any(text.startswith("- ") for text in texts),
                        f"{path.name}:slide_{slide_index}",
                    )


if __name__ == "__main__":
    unittest.main()
