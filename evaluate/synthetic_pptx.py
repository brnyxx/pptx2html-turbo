from __future__ import annotations

import argparse
from collections.abc import Sequence
from pathlib import Path
from typing import cast
from zipfile import ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu as PptxEmu

from evaluate.synthetic_scene import (
    CANVAS_HEIGHT_EMU,
    CANVAS_WIDTH_EMU,
    RgbHex,
    SyntheticDeck,
    create_synthetic_corpus,
)


def write_synthetic_pptx_corpus(
    corpus: Sequence[SyntheticDeck],
    output_dir: Path,
) -> tuple[Path, ...]:
    output_dir.mkdir(parents=True, exist_ok=True)
    outputs: list[Path] = []
    for deck in corpus:
        presentation = Presentation()
        presentation.slide_width = PptxEmu(CANVAS_WIDTH_EMU)
        presentation.slide_height = PptxEmu(CANVAS_HEIGHT_EMU)
        for scene in deck.scenes:
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            background = slide.background.fill
            background.solid()
            background.fore_color.rgb = _rgb(scene.background)
            for rectangle in scene.rectangles:
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    PptxEmu(rectangle.x),
                    PptxEmu(rectangle.y),
                    PptxEmu(rectangle.width),
                    PptxEmu(rectangle.height),
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = _rgb(rectangle.fill)
                shape.line.fill.background()
        output = output_dir / f"{deck.name}.pptx"
        presentation.save(str(output))
        _disable_theme_effects(output)
        outputs.append(output)
    return tuple(outputs)


def main(argv: Sequence[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        description="Generate the deterministic synthetic oracle deck corpus."
    )
    _ = parser.add_argument("--output", type=Path, required=True)
    args = parser.parse_args(argv)

    output = cast(Path, args.output)
    _ = write_synthetic_pptx_corpus(create_synthetic_corpus(), output)
    return 0


def _rgb(color: RgbHex) -> RGBColor:
    return RGBColor(
        int(color[0:2], 16),
        int(color[2:4], 16),
        int(color[4:6], 16),
    )


def _disable_theme_effects(path: Path) -> None:
    temporary = path.with_suffix(".pptx.tmp")
    temporary.unlink(missing_ok=True)
    replacements = 0
    try:
        with ZipFile(path) as source, ZipFile(temporary, "w") as target:
            for part in source.infolist():
                payload = source.read(part)
                if part.filename.startswith(
                    "ppt/slides/slide"
                ) and part.filename.endswith(".xml"):
                    replacements += payload.count(b'<a:effectRef idx="2">')
                    payload = payload.replace(
                        b'<a:effectRef idx="2">',
                        b'<a:effectRef idx="0">',
                    )
                target.writestr(part, payload)
        if replacements == 0:
            raise ValueError(f"SYNTHETIC_EFFECT_REFERENCE_MISSING:{path}")
        with ZipFile(temporary) as validation:
            corrupt_part = validation.testzip()
        if corrupt_part is not None:
            raise ValueError(
                f"SYNTHETIC_PACKAGE_INVALID:{path}:{corrupt_part}"
            )
        _ = temporary.replace(path)
    finally:
        temporary.unlink(missing_ok=True)


if __name__ == "__main__":
    raise SystemExit(main())
